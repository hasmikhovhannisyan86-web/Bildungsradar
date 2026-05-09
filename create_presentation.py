"""BildungsRadar Präsentation erstellen."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Farben
NAVY = RGBColor(0x1E, 0x27, 0x61)
ICE_BLUE = RGBColor(0xCA, 0xDC, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT = RGBColor(0x3D, 0x5A, 0xF1)
LIGHT_BG = RGBColor(0xF4, 0xF6, 0xFC)
MEDIUM_BLUE = RGBColor(0x2D, 0x3A, 0x8C)
GREEN_ACCENT = RGBColor(0x27, 0xAE, 0x60)
ORANGE_ACCENT = RGBColor(0xF5, 0xA6, 0x23)
SUBTLE_GRAY = RGBColor(0x8E, 0x8E, 0xA3)


def add_shape(slide, left, top, width, height, fill_color, border=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if not border:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=16,
                 color=DARK_TEXT, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=14,
                    color=DARK_TEXT, icon="▸"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{icon} {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(6)
    return txBox


def add_icon_circle(slide, left, top, size, color, text, text_size=18):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(text_size)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)


# ============================================================
# SLIDE 1 - Titelfolie
# ============================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

# Hintergrund
bg = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = NAVY
bg.line.fill.background()

# Linker Akzent-Streifen
stripe = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.15), prs.slide_height)
stripe.fill.solid()
stripe.fill.fore_color.rgb = ACCENT
stripe.line.fill.background()

# Dezente Linie
line_shape = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2), Inches(3.7), Inches(9.3), Pt(2))
line_shape.fill.solid()
line_shape.fill.fore_color.rgb = ACCENT
line_shape.line.fill.background()

add_text_box(slide1, Inches(2), Inches(1.5), Inches(9.3), Inches(1.5),
             "BildungsRadar", font_size=54, color=WHITE, bold=True, font_name="Arial Black")

add_text_box(slide1, Inches(2), Inches(2.8), Inches(9.3), Inches(0.8),
             "KI-gestuetzte Suche und Analyse von Bildungseinrichtungen",
             font_size=22, color=ICE_BLUE, font_name="Calibri")

add_text_box(slide1, Inches(2), Inches(4.2), Inches(9.3), Inches(0.5),
             "Hasmik Hovhannisyan", font_size=20, color=WHITE, bold=True)

add_text_box(slide1, Inches(2), Inches(4.8), Inches(9.3), Inches(0.5),
             "AI Engineering  |  Masterschool  |  2026", font_size=16, color=SUBTLE_GRAY)

# GitHub Link
add_text_box(slide1, Inches(2), Inches(5.8), Inches(9.3), Inches(0.4),
             "github.com/hasmikhovhannisyan86-web/Bildungsradar",
             font_size=13, color=SUBTLE_GRAY)


# ============================================================
# SLIDE 2 - Problemstellung & Lösung
# ============================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])

# Heller Hintergrund
bg2 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg2.fill.solid()
bg2.fill.fore_color.rgb = LIGHT_BG
bg2.line.fill.background()

# Header-Balken
header2 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
header2.fill.solid()
header2.fill.fore_color.rgb = NAVY
header2.line.fill.background()

add_text_box(slide2, Inches(0.8), Inches(0.25), Inches(11), Inches(0.7),
             "Problemstellung & Loesung", font_size=36, color=WHITE, bold=True, font_name="Arial Black")

# Problem-Karte (links)
prob_card = add_shape(slide2, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.8), WHITE)

add_icon_circle(slide2, Inches(1.3), Inches(2.2), Inches(0.7), RGBColor(0xE7, 0x4C, 0x3C), "?", 24)

add_text_box(slide2, Inches(2.2), Inches(2.25), Inches(3.5), Inches(0.5),
             "Das Problem", font_size=24, color=RGBColor(0xE7, 0x4C, 0x3C), bold=True)

prob_items = [
    "Eltern suchen stundenlang nach der richtigen Schule oder Kita",
    "Informationen verstreut ueber verschiedene Webseiten",
    "Kein einfacher Vergleich zwischen Einrichtungen moeglich",
    "Preise, Angebote und Bewertungen schwer zu finden"
]
add_bullet_list(slide2, Inches(1.3), Inches(3.2), Inches(4.5), Inches(3.0),
                prob_items, font_size=14, color=DARK_TEXT)

# Lösung-Karte (rechts)
sol_card = add_shape(slide2, Inches(7.0), Inches(1.8), Inches(5.5), Inches(4.8), WHITE)

add_icon_circle(slide2, Inches(7.5), Inches(2.2), Inches(0.7), GREEN_ACCENT, "!", 24)

add_text_box(slide2, Inches(8.4), Inches(2.25), Inches(3.5), Inches(0.5),
             "Die Loesung", font_size=24, color=GREEN_ACCENT, bold=True)

sol_items = [
    "BildungsRadar: Eine Web-App fuer Eltern",
    "Automatische Suche nach Ort (OpenStreetMap)",
    "KI analysiert Webseiten der Einrichtungen",
    "Vergleichstabelle: Alle Infos auf einen Blick"
]
add_bullet_list(slide2, Inches(7.5), Inches(3.2), Inches(4.5), Inches(3.0),
                sol_items, font_size=14, color=DARK_TEXT)


# ============================================================
# SLIDE 3 - Technische Architektur
# ============================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])

bg3 = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg3.fill.solid()
bg3.fill.fore_color.rgb = LIGHT_BG
bg3.line.fill.background()

header3 = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
header3.fill.solid()
header3.fill.fore_color.rgb = NAVY
header3.line.fill.background()

add_text_box(slide3, Inches(0.8), Inches(0.25), Inches(11), Inches(0.7),
             "Technische Architektur", font_size=36, color=WHITE, bold=True, font_name="Arial Black")

# 2x3 Grid von Tech-Karten
cards = [
    ("Frontend", "Flask + HTML/CSS/JS", ACCENT, "W"),
    ("Datenquelle", "OpenStreetMap\nOverpass API", MEDIUM_BLUE, "M"),
    ("KI-Analyse", "OpenAI GPT-4o-mini\n+ Web-Scraping", RGBColor(0x9C, 0x27, 0xB0), "K"),
    ("Datenbank", "SQLite\nfuer Caching", GREEN_ACCENT, "D"),
    ("Filter", "Kindergaerten, Kitas\nSchulen, Privatschulen", ORANGE_ACCENT, "F"),
    ("Vergleich", "Einrichtungen\nnebeneinander", RGBColor(0xE7, 0x4C, 0x3C), "V"),
]

for i, (title, desc, color, icon) in enumerate(cards):
    col = i % 3
    row = i // 3
    x = Inches(0.8 + col * 4.1)
    y = Inches(1.7 + row * 2.7)

    card = add_shape(slide3, x, y, Inches(3.7), Inches(2.3), WHITE)

    add_icon_circle(slide3, x + Inches(0.3), y + Inches(0.3), Inches(0.6), color, icon, 18)

    add_text_box(slide3, x + Inches(1.1), y + Inches(0.3), Inches(2.2), Inches(0.4),
                 title, font_size=18, color=DARK_TEXT, bold=True)

    add_text_box(slide3, x + Inches(0.3), y + Inches(1.1), Inches(3.1), Inches(1.0),
                 desc, font_size=14, color=SUBTLE_GRAY)

# Stat-Callout
add_text_box(slide3, Inches(9.0), Inches(6.2), Inches(3.5), Inches(0.8),
             "935+ Einrichtungen pro Stadt", font_size=16, color=ACCENT, bold=True,
             alignment=PP_ALIGN.RIGHT)


# ============================================================
# SLIDE 4 - Prompt Engineering Ergebnisse
# ============================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])

bg4 = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg4.fill.solid()
bg4.fill.fore_color.rgb = LIGHT_BG
bg4.line.fill.background()

header4 = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
header4.fill.solid()
header4.fill.fore_color.rgb = NAVY
header4.line.fill.background()

add_text_box(slide4, Inches(0.8), Inches(0.25), Inches(11), Inches(0.7),
             "Prompt Engineering: 3 Varianten im Vergleich", font_size=34, color=WHITE,
             bold=True, font_name="Arial Black")

# Drei Prompt-Karten nebeneinander
prompt_cards = [
    ("v1 Basic", "Einfache direkte\nAnweisung an die KI",
     RGBColor(0x4C, 0xAF, 0x50), ["Schnell", "Weniger Details", "Preise oft fehlend"]),
    ("v2 Few-Shot", "Prompt mit Beispiel\nzur Orientierung",
     RGBColor(0x21, 0x96, 0xF3), ["Mittlere Details", "Bewertungen gefunden", "4.5 Sterne (Phorms)"]),
    ("v3 Chain-of-Thought", "Schritt-fuer-Schritt\nAnalyse-Anleitung",
     RGBColor(0x9C, 0x27, 0xB0), ["Detaillierteste Ergebnisse", "Preise erkannt", "Beste Qualitaet"]),
]

for i, (title, desc, color, results) in enumerate(prompt_cards):
    x = Inches(0.8 + i * 4.1)
    y = Inches(1.7)

    card = add_shape(slide4, x, y, Inches(3.7), Inches(5.0), WHITE)

    # Farbiger Header im Card
    card_header = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.7), Inches(0.8))
    card_header.fill.solid()
    card_header.fill.fore_color.rgb = color
    card_header.line.fill.background()

    add_text_box(slide4, x + Inches(0.3), y + Inches(0.15), Inches(3.1), Inches(0.5),
                 title, font_size=20, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide4, x + Inches(0.3), y + Inches(1.1), Inches(3.1), Inches(0.8),
                 desc, font_size=14, color=SUBTLE_GRAY, alignment=PP_ALIGN.CENTER)

    add_text_box(slide4, x + Inches(0.3), y + Inches(2.1), Inches(3.1), Inches(0.4),
                 "Ergebnisse:", font_size=14, color=DARK_TEXT, bold=True)

    add_bullet_list(slide4, x + Inches(0.3), y + Inches(2.6), Inches(3.1), Inches(2.0),
                    results, font_size=13, color=DARK_TEXT, icon="✓")

# Winner-Badge
winner = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(9.2), Inches(5.6), Inches(3.3), Inches(0.5))
winner.fill.solid()
winner.fill.fore_color.rgb = RGBColor(0x9C, 0x27, 0xB0)
winner.line.fill.background()
tf = winner.text_frame
p = tf.paragraphs[0]
p.text = "★  v3 = Beste Ergebnisse"
p.font.size = Pt(14)
p.font.color.rgb = WHITE
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

# Temperature info
add_text_box(slide4, Inches(0.8), Inches(6.8), Inches(6), Inches(0.4),
             "Optimale Temperature: 0.3 (praezise Faktenextraktion)",
             font_size=13, color=SUBTLE_GRAY)


# ============================================================
# SLIDE 5 - Demo & Fazit
# ============================================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])

# Dunkler Hintergrund für Abschluss
bg5 = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg5.fill.solid()
bg5.fill.fore_color.rgb = NAVY
bg5.line.fill.background()

stripe5 = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.15), prs.slide_height)
stripe5.fill.solid()
stripe5.fill.fore_color.rgb = ACCENT
stripe5.line.fill.background()

add_text_box(slide5, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "Demo & Fazit", font_size=40, color=WHITE, bold=True, font_name="Arial Black")

# Stat-Callouts (große Zahlen)
stats = [
    ("935+", "Einrichtungen\npro Stadt"),
    ("37", "Privatschulen\nin Frankfurt"),
    ("3", "Prompt-Varianten\ngetestet"),
]

for i, (number, label) in enumerate(stats):
    x = Inches(0.8 + i * 4.1)
    y = Inches(1.8)

    stat_card = add_shape(slide5, x, y, Inches(3.7), Inches(1.8), MEDIUM_BLUE)

    add_text_box(slide5, x + Inches(0.3), y + Inches(0.15), Inches(3.1), Inches(0.9),
                 number, font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER,
                 font_name="Arial Black")

    add_text_box(slide5, x + Inches(0.3), y + Inches(1.1), Inches(3.1), Inches(0.6),
                 label, font_size=14, color=ICE_BLUE, alignment=PP_ALIGN.CENTER)

# Features & Fazit
features = [
    "Live-Suche nach Bildungseinrichtungen in jeder deutschen Stadt",
    "KI-Analyse extrahiert Angebote, Preise und Spezialisierungen",
    "Vergleichstabelle: Schulen nebeneinander vergleichen",
    "Prompt-Vergleich: v1 vs v2 vs v3 auf einen Blick sichtbar",
]
add_bullet_list(slide5, Inches(0.8), Inches(4.1), Inches(7), Inches(2.5),
                features, font_size=15, color=ICE_BLUE, icon="→")

# Fazit-Box
fazit = add_shape(slide5, Inches(8.5), Inches(4.1), Inches(4.2), Inches(2.0), ACCENT)
add_text_box(slide5, Inches(8.8), Inches(4.3), Inches(3.6), Inches(0.4),
             "Fazit", font_size=20, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide5, Inches(8.8), Inches(4.85), Inches(3.6), Inches(1.0),
             "KI kann Eltern bei der Schulwahl unterstuetzen und spart wertvolle Zeit.",
             font_size=14, color=WHITE, alignment=PP_ALIGN.CENTER)

# GitHub
add_text_box(slide5, Inches(0.8), Inches(6.7), Inches(11), Inches(0.4),
             "github.com/hasmikhovhannisyan86-web/Bildungsradar",
             font_size=14, color=SUBTLE_GRAY)

# Speichern
output_path = "/Users/hasmikhovhannisyan/Desktop/Projektarbeit/BildungsRadar_Praesentation.pptx"
prs.save(output_path)
print(f"Praesentation gespeichert: {output_path}")
